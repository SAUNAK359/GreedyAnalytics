from pptx import Presentation


def load_ppt(file) -> dict:
    prs = Presentation(file)
    text_runs = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_runs.append(shape.text)
    return {"type": "document", "text": "\n".join(text_runs)}
